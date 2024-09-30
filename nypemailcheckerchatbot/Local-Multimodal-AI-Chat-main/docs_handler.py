from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from vectordb_handler import load_vectordb
from utils import load_config, timeit
import io
import pypdfium2
import docx
import openpyxl
from pptx import Presentation

config = load_config()

def get_pdf_texts(pdfs_bytes_list):
    return [extract_text_from_pdf(io.BytesIO(pdf_bytes.read())) for pdf_bytes in pdfs_bytes_list]

def extract_text_from_pdf(pdf_bytes_io):
    pdf_file = pypdfium2.PdfDocument(pdf_bytes_io)
    return "\n".join(pdf_file.get_page(page_number).get_textpage().get_text_range() for page_number in range(len(pdf_file)))

def get_docx_texts(docx_bytes_list):
    return [extract_text_from_docx(io.BytesIO(docx_bytes.read())) for docx_bytes in docx_bytes_list]

def extract_text_from_docx(docx_bytes_io):
    doc = docx.Document(docx_bytes_io)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)

def get_xlsx_texts(xlsx_bytes_list):
    return [extract_text_from_xlsx(io.BytesIO(xlsx_bytes.read())) for xlsx_bytes in xlsx_bytes_list]

def extract_text_from_xlsx(xlsx_bytes_io):
    workbook = openpyxl.load_workbook(xlsx_bytes_io)
    text = []
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
            text.append(row_text)
    return "\n".join(text)

def get_pptx_texts(pptx_bytes_list):
    return [extract_text_from_pptx(io.BytesIO(pptx_bytes.read())) for pptx_bytes in pptx_bytes_list]

def extract_text_from_pptx(pptx_bytes_io):
    prs = Presentation(pptx_bytes_io)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=config["text_splitter"]["chunk_size"], 
                                              chunk_overlap=config["text_splitter"]["overlap"],
                                              separators=config["text_splitter"]["separators"])
    return splitter.split_text(text)

def get_document_chunks(text_list):
    documents = []
    for text in text_list:
        for chunk in get_text_chunks(text):
            documents.append(Document(page_content=chunk))
    return documents

@timeit
def add_documents_to_db(files_bytes, file_type):
    if file_type == "pdf":
        texts = get_pdf_texts(files_bytes)
    elif file_type == "docx":
        texts = get_docx_texts(files_bytes)
    elif file_type == "xlsx":
        texts = get_xlsx_texts(files_bytes)
    elif file_type == "pptx":
        texts = get_pptx_texts(files_bytes)
    else:
        raise ValueError("Unsupported file type")

    documents = get_document_chunks(texts)
    vector_db = load_vectordb()
    vector_db.add_documents(documents)
    print(f"{file_type.upper()} documents added to db.")